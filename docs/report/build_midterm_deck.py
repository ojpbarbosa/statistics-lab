"""Generate the midterm check-in PowerPoint decks (EN + PT).

Run:  python3 docs/report/build_midterm_deck.py
Writes midterm_presentation_en.pptx / midterm_presentation_pt.pptx next to
this file. Pure python-pptx, 16:9, Segoe UI, no bullet points anywhere.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent

FONT = "Segoe UI"
NAVY = RGBColor(0x16, 0x38, 0x5C)
NAVY_SOFT = RGBColor(0x3D, 0x5F, 0x8A)
INK = RGBColor(0x1F, 0x29, 0x37)
TINT = RGBColor(0xE9, 0xEE, 0xF5)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
FOOT = RGBColor(0x5A, 0x70, 0x89)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)
WARN = RGBColor(0xFA, 0xB2, 0x19)
CRIT = RGBColor(0xD0, 0x3B, 0x3B)
SLATE = RGBColor(0x4B, 0x55, 0x63)
# dataviz-validated categorical slots for the weights bar
W_COLORS = [RGBColor(0x2A, 0x78, 0xD6), RGBColor(0x00, 0x83, 0x00),
            RGBColor(0xE8, 0x7B, 0xA4), RGBColor(0xED, 0xA1, 0x00),
            RGBColor(0x1B, 0xAF, 0x7A)]
W_DARK_TEXT = [False, False, True, True, True]

ML, SLIDE_W, SLIDE_H = 0.75, 13.333, 7.5
CW = SLIDE_W - 2 * ML  # content width 11.83


def para(tf, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
         space_after=6, italic=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def box(slide, x, y, w, h, fill, rounded=True, radius=0.10):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    if rounded:
        try:
            shape.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def card(slide, x, y, w, h, title, body, title_size=15, body_size=11.5):
    shape = box(slide, x, y, w, h, TINT)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, title, size=title_size, bold=True, color=NAVY, space_after=4, first=True)
    para(tf, body, size=body_size, color=INK, space_after=0)
    return shape


def pill(slide, x, y, w, label, fill, dark_text=False):
    shape = box(slide, x, y, w, 0.34, fill, radius=0.5)
    para(shape.text_frame, label, size=10.5, bold=True,
         color=INK if dark_text else PAPER, align=PP_ALIGN.CENTER, space_after=0, first=True)
    return shape


def arrow(slide, x, y):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.38), Inches(0.26))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_SOFT
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def flow_node(slide, x, y, w, h, title, sub, fill=TINT, title_color=NAVY):
    shape = box(slide, x, y, w, h, fill)
    tf = shape.text_frame
    para(tf, title, size=13, bold=True, color=title_color, align=PP_ALIGN.CENTER, space_after=2, first=True)
    para(tf, sub, size=9.5, color=INK, align=PP_ALIGN.CENTER, space_after=0)
    return shape


def content_slide(prs, title, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, ML, 0.55, 0.55, 0.09, NAVY, rounded=False)
    tf = textbox(slide, ML, 0.72, CW, 0.7)
    para(tf, title, size=26, bold=True, color=INK, space_after=0, first=True)
    foot = textbox(slide, ML, SLIDE_H - 0.42, CW, 0.3)
    p = para(foot, footer, size=9, color=FOOT, space_after=0, first=True)
    p.alignment = PP_ALIGN.LEFT
    return slide


def intro(slide, text, y=1.55, size=13, h=0.75):
    tf = textbox(slide, ML, y, CW, h)
    para(tf, text, size=size, color=INK, space_after=0, first=True)


STR = {
    "en": {
        "footer": "Issuer Opportunity Screener · midterm check-in",
        "title": "Issuer Opportunity Screener",
        "subtitle": "Midterm check-in",
        "byline": "João Pedro Ferreira Barbosa",
        "date": "July 21, 2026 · week 2 of 4",
        "tagline": "A screening framework that looks for corporate names whose spread over the Brazil benchmark makes an offshore note worth pitching to Brazilian investors.",

        "s2_title": "What I want from this conversation",
        "s2_intro": "Quick framing before anything else: this is not the final readout, it is a midterm alignment. I want to show you where things stand, sanity-check a few decisions while they are still cheap to change, and set up a faster feedback loop for the remaining weeks.",
        "s2_cards": [
            ("Report", "Walk through what is built, what already works, and what the data looks like today."),
            ("Validate", "Check the framework and criteria choices with you: the rule, the weights, the tolerances."),
            ("Unblock", "Be transparent about where data access stands and the routes I am working around it."),
            ("Align", "Agree on the next steps and on checking in more often, in smaller doses."),
        ],

        "s3_title": "Where the project stands",
        "s3_intro": "The groundwork phase is essentially done. The screener runs end to end today, with synthetic data standing in wherever live access is still pending.",
        "s3_tiles": [("125", "names in the screening universe"),
                     ("100+", "automated tests, runs on any machine"),
                     ("1", "open blocker, and it is external")],
        "s3_plan": "The four week plan",
        "s3_weeks": ["Week 1 · done", "Week 2 · in progress", "Week 3 · results", "Week 4 · delivery"],
        "s3_today": "today",

        "s4_title": "How it is put together",
        "s4_intro": "One pipeline, five stages, each replaceable without touching the others. That is what lets me swap data sources while everything downstream stays identical.",
        "s4_nodes": [
            ("Universe", "125 names, add, quarantine, restore"),
            ("Sources", "BLP, BQuant, Hermes, fixture"),
            ("Snapshots", "versioned and append-only"),
            ("Scoring", "5 blocks, tiers A, B, C"),
            ("Dashboard", "screen, movers, reports"),
        ],
        "s4_note": "Every classification carries a plain-language rationale, and every number can be replicated on the terminal.",

        "s5_title": "Where data access stands",
        "s5_intro": "Four routes to the data, in different states. Whichever lands first feeds the exact same pipeline, so no work gets thrown away.",
        "s5_rows": [
            ("Bloomberg Desktop API", "Bulk bond requests are gated behind a workflow review. I contacted our account manager and I am still waiting to hear back.", "blocked", CRIT, False),
            ("BQuant", "The server-side bond screen is integrated and working. I still owe it a proper validation run on the Terminal.", "validating", WARN, True),
            ("Hermes internal API", "Found this week. Bonds already flow through the fetch endpoint, and I am talking to the owner about exposing CDS server side.", "bonds live", GOOD, False),
            ("Markit Partners", "Access requested, for additional credit data to enrich the screen.", "requested", SLATE, False),
        ],

        "s6_title": "How the screening decides",
        "s6_intro": "The core rule reads in one line: a name is worth a look when it pays at least what Brazil pays. Names slightly inside Brazil survive only with a clearly stronger rating, which is how Hyundai at minus 18 bps and A minus stays in.",
        "s6_zones": ["not viable", "edge case", "viable"],
        "s6_ticks": ["-20 bps", "0"],
        "s6_weights_label": "What the composite score weighs",
        "s6_weights": ["Spread", "Credit", "Liquidity", "Equity", "Recognition"],
        "s6_note": "The order mirrors what the desk cares about: carry first, then credit safety, tradability, the market's read, and how recognizable the name is to clients. Scores land in tiers: A is pitch ready, B is a watchlist, C is a pass.",

        "s7_title": "A few things I would like your read on",
        "s7_intro": "None of these block me, but your answers change what I build next.",
        "s7_cards": [
            ("Criteria", "Does the 20 bps tolerance feel right? Should Sr Non-Preferred paper count as senior? Any currency preference beyond USD?"),
            ("Universe file", "I need internal ratings, ISINs for the Hermes route, and handle overrides for the non-US names."),
            ("Dashboard", "Is any visualization or metric missing for how you would actually use this day to day?"),
            ("Cadence", "Would a short weekly check-in work for you, instead of one big handoff at the end?"),
        ],
        "s7_note": "One ask on the feedback itself: reactions to the framework and criteria move the project further than reactions to individual names.",

        "s8_title": "A challenge I want to bring: growing the universe",
        "s8_intro": "Today every name enters the universe by hand, one form at a time, and the list stays frozen until someone edits it. From a product point of view that bothers me: we only ever analyze the names we already knew about.",
        "s8_now_label": "How it works today",
        "s8_now": [("A name is added by hand", "one form per name"), ("Static universe", "the same 125 names until someone edits")],
        "s8_prop_label": "What I would explore",
        "s8_prop": [
            ("Feeds", "Hermes tape, index members, new issues"),
            ("Auto pre-screen", "spread vs Brazil, liquidity"),
            ("Candidate inbox", "ranked suggestions"),
            ("Desk approves", "the final say stays human"),
        ],
        "s8_note": "Hermes already returns the whole bond tape for any date. Diffing that tape against our universe gives machine-suggested candidates with spreads precomputed. Not manual, not brute force, and the desk keeps the last word.",

        "s9_title": "The same lens, applied elsewhere",
        "s9_intro": "This is where I think I add the most value: taking an engineering and product view of the desk's workflow and removing manual steps between a signal and a decision. Three places I would apply it next:",
        "s9_cards": [
            ("Built-in credit analysis", "An automatic one-pager per name: balance sheet, leverage, covenants and news, assembled by the system so the analyst validates instead of gathering."),
            ("From screen to monitor", "Scheduled snapshots plus alerts on viability flips and unusual moves, so the desk gets told instead of having to look."),
            ("A learning loop", "Your feedback on names, plus backtests over the accumulated snapshots, becomes evidence to tune the score weights."),
        ],

        "s10_title": "Next steps, and how I would like to work",
        "s10_cols": [
            ("This week, Wed to Fri", "Validate the Hermes route end to end and land the first complete snapshot of the full universe."),
            ("Week 3", "Results pass: coverage, tier distribution, top names and movers. Plus your pass on the universe file."),
            ("Week 4", "Categorized candidate lists, the final report, and the closing presentation with the code handover."),
        ],
        "s10_cadence": "And instead of one big handoff at the end, I would rather do a short check-in each week, so we can course-correct while it is still cheap.",
        "s10_chip": "check-in",
    },
    "pt": {
        "footer": "Issuer Opportunity Screener · check-in de meio de percurso",
        "title": "Issuer Opportunity Screener",
        "subtitle": "Check-in de meio de percurso",
        "byline": "João Pedro Ferreira Barbosa",
        "date": "21/07/2026 · semana 2 de 4",
        "tagline": "Um framework de screening que procura nomes corporativos cujo spread sobre o benchmark Brasil torna uma nota offshore digna de pitch para investidores brasileiros.",

        "s2_title": "O que eu quero desta conversa",
        "s2_intro": "Um enquadramento rápido antes de tudo: esta não é a entrega final, é um alinhamento de meio de percurso. Quero mostrar onde as coisas estão, validar algumas decisões enquanto ainda é barato mudá-las e montar um ciclo de feedback mais rápido para as semanas restantes.",
        "s2_cards": [
            ("Reportar", "Passar pelo que está construído, o que já funciona e como estão os dados hoje."),
            ("Validar", "Checar com vocês as escolhas de framework e critérios: a regra, os pesos, as tolerâncias."),
            ("Destravar", "Ser transparente sobre o acesso a dados e as rotas que estou construindo em volta."),
            ("Alinhar", "Combinar os próximos passos e conversar mais vezes, em doses menores."),
        ],

        "s3_title": "Onde o projeto está",
        "s3_intro": "A fase de fundação está essencialmente pronta. O screener roda de ponta a ponta hoje, com dados sintéticos cobrindo onde o acesso live ainda está pendente.",
        "s3_tiles": [("125", "nomes no universo de screening"),
                     ("100+", "testes automatizados, roda em qualquer máquina"),
                     ("1", "bloqueio aberto, e ele é externo")],
        "s3_plan": "O plano de quatro semanas",
        "s3_weeks": ["Semana 1 · feita", "Semana 2 · em andamento", "Semana 3 · resultados", "Semana 4 · entrega"],
        "s3_today": "hoje",

        "s4_title": "Como está montado",
        "s4_intro": "Um pipeline, cinco estágios, cada um substituível sem tocar nos outros. É isso que me deixa trocar fontes de dados com tudo que vem depois permanecendo idêntico.",
        "s4_nodes": [
            ("Universo", "125 nomes, incluir, quarentenar, restaurar"),
            ("Fontes", "BLP, BQuant, Hermes, fixture"),
            ("Snapshots", "versionados e append-only"),
            ("Scoring", "5 blocos, tiers A, B, C"),
            ("Dashboard", "screen, movers, relatórios"),
        ],
        "s4_note": "Toda classificação carrega uma justificativa em linguagem simples, e todo número pode ser replicado no terminal.",

        "s5_title": "Onde está o acesso a dados",
        "s5_intro": "Quatro rotas até os dados, em estados diferentes. A que chegar primeiro alimenta exatamente o mesmo pipeline, então nenhum trabalho é jogado fora.",
        "s5_rows": [
            ("Bloomberg Desktop API", "Requisições de bonds em lote estão travadas atrás de uma revisão de workflow. Contatei nosso account manager e ainda estou aguardando retorno.", "bloqueada", CRIT, False),
            ("BQuant", "O screening de bonds no servidor está integrado e funcionando. Ainda devo uma rodada de validação no Terminal.", "validando", WARN, True),
            ("API interna Hermes", "Descoberta esta semana. Bonds já fluem pelo endpoint de fetch, e estou conversando com o responsável sobre expor CDS no servidor.", "bonds no ar", GOOD, False),
            ("Markit Partners", "Acesso solicitado, para dados adicionais de crédito enriquecerem o screen.", "solicitada", SLATE, False),
        ],

        "s6_title": "Como o screening decide",
        "s6_intro": "A regra central cabe em uma linha: um nome merece atenção quando paga pelo menos o que o Brasil paga. Nomes um pouco abaixo do Brasil só sobrevivem com rating claramente mais forte, e é assim que a Hyundai, a menos 18 bps com A menos, continua dentro.",
        "s6_zones": ["não viável", "edge case", "viável"],
        "s6_ticks": ["-20 bps", "0"],
        "s6_weights_label": "O que o score composto pesa",
        "s6_weights": ["Spread", "Crédito", "Liquidez", "Equity", "Reconhecimento"],
        "s6_note": "A ordem espelha o que importa para a mesa: carrego primeiro, depois segurança de crédito, negociabilidade, a leitura do mercado e o quanto o nome é reconhecível para o cliente. Os scores caem em tiers: A pronto para pitch, B watchlist, C descarte.",

        "s7_title": "Algumas leituras que eu queria de vocês",
        "s7_intro": "Nenhuma delas me bloqueia, mas as respostas mudam o que eu construo em seguida.",
        "s7_cards": [
            ("Critérios", "A tolerância de 20 bps parece certa? Papel Sr Non-Preferred conta como sênior? Alguma preferência de moeda além de USD?"),
            ("Arquivo do universo", "Preciso de ratings internos, ISINs para a rota Hermes e overrides de handles para os nomes fora dos EUA."),
            ("Dashboard", "Falta alguma visualização ou métrica para o uso real de vocês no dia a dia?"),
            ("Cadência", "Um check-in curto por semana funciona para vocês, em vez de uma única grande entrega no final?"),
        ],
        "s7_note": "Um pedido sobre o próprio feedback: reações ao framework e aos critérios movem o projeto mais do que reações a nomes individuais.",

        "s8_title": "Um desafio que quero trazer: crescer o universo",
        "s8_intro": "Hoje todo nome entra no universo à mão, um formulário por vez, e a lista fica congelada até alguém editar. Do ponto de vista de produto isso me incomoda: só analisamos os nomes que já conhecíamos.",
        "s8_now_label": "Como funciona hoje",
        "s8_now": [("Nome incluído à mão", "um formulário por nome"), ("Universo estático", "os mesmos 125 nomes até alguém editar")],
        "s8_prop_label": "O que eu exploraria",
        "s8_prop": [
            ("Feeds", "fita do Hermes, membros de índices, novas emissões"),
            ("Pré-screen automático", "spread vs Brasil, liquidez"),
            ("Inbox de candidatos", "sugestões ranqueadas"),
            ("Mesa aprova", "a palavra final continua humana"),
        ],
        "s8_note": "O Hermes já retorna a fita inteira de bonds para qualquer data. Comparar essa fita com o nosso universo gera candidatos sugeridos pela máquina com spreads já calculados. Nem manual, nem força bruta, e a mesa mantém a última palavra.",

        "s9_title": "A mesma lente, aplicada em outros lugares",
        "s9_intro": "É aqui que eu acho que agrego mais: olhar o fluxo da mesa com uma visão de engenharia e produto e remover passos manuais entre um sinal e uma decisão. Três lugares onde eu aplicaria isso em seguida:",
        "s9_cards": [
            ("Análise de crédito embutida", "Um one-pager automático por nome: balanço, alavancagem, covenants e notícias, montado pelo sistema para o analista validar em vez de coletar."),
            ("De screen a monitor", "Snapshots agendados mais alertas de viradas de viabilidade e movimentos incomuns, para a mesa ser avisada em vez de precisar olhar."),
            ("Um loop de aprendizado", "O feedback de vocês sobre nomes, mais backtests sobre os snapshots acumulados, vira evidência para calibrar os pesos do score."),
        ],

        "s10_title": "Próximos passos, e como eu gostaria de trabalhar",
        "s10_cols": [
            ("Esta semana, qua a sex", "Validar a rota Hermes de ponta a ponta e gravar o primeiro snapshot completo do universo inteiro."),
            ("Semana 3", "Rodada de resultados: cobertura, distribuição de tiers, principais nomes e movers. Mais a rodada de vocês no arquivo do universo."),
            ("Semana 4", "Listas categorizadas de candidatos, o relatório final e a apresentação de encerramento com o handover do código."),
        ],
        "s10_cadence": "E em vez de uma única grande entrega no final, eu preferia um check-in curto por semana, para corrigirmos o rumo enquanto ainda é barato.",
        "s10_chip": "check-in",
    },
}


def build(lang: str) -> Path:
    L = STR[lang]
    prs = Presentation()
    prs.slide_width = Emu(int(Inches(SLIDE_W)))
    prs.slide_height = Emu(int(Inches(SLIDE_H)))

    # ---- 1 · title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, 5.4, SLIDE_H, NAVY, rounded=False)
    tf = textbox(slide, 0.6, 2.3, 4.4, 3.4)
    para(tf, L["title"], size=32, bold=True, color=PAPER, space_after=8, first=True)
    para(tf, L["subtitle"], size=17, color=TINT, space_after=26)
    para(tf, L["byline"], size=12, bold=True, color=PAPER, space_after=2)
    para(tf, L["date"], size=11, color=TINT, space_after=0)
    tf = textbox(slide, 6.1, 3.0, 6.4, 1.8)
    para(tf, L["tagline"], size=15, color=INK, space_after=0, first=True)
    box(slide, 6.1, 2.7, 0.55, 0.09, NAVY, rounded=False)

    # ---- 2 · purpose ----
    slide = content_slide(prs, L["s2_title"], L["footer"])
    intro(slide, L["s2_intro"], h=0.95)
    y, w, gap = 2.75, 2.83, 0.17
    for i, (title, body) in enumerate(L["s2_cards"]):
        card(slide, ML + i * (w + gap), y, w, 2.35, title, body)

    # ---- 3 · status ----
    slide = content_slide(prs, L["s3_title"], L["footer"])
    intro(slide, L["s3_intro"])
    y, w, gap = 2.35, 3.83, 0.17
    for i, (num, label) in enumerate(L["s3_tiles"]):
        shape = box(slide, ML + i * (w + gap), y, w, 1.75, NAVY)
        tf = shape.text_frame
        para(tf, num, size=38, bold=True, color=PAPER, align=PP_ALIGN.CENTER, space_after=2, first=True)
        para(tf, label, size=11.5, color=TINT, align=PP_ALIGN.CENTER, space_after=0)
    tf = textbox(slide, ML, 4.65, CW, 0.35)
    para(tf, L["s3_plan"], size=12, bold=True, color=INK, space_after=0, first=True)
    y, w, gap = 5.05, 2.925, 0.043
    fills = [NAVY, NAVY_SOFT, TINT, TINT]
    texts = [PAPER, PAPER, NAVY, NAVY]
    for i, label in enumerate(L["s3_weeks"]):
        shape = box(slide, ML + i * (w + gap), y, w, 0.55, fills[i], radius=0.18)
        para(shape.text_frame, label, size=11.5, bold=i < 2, color=texts[i],
             align=PP_ALIGN.CENTER, space_after=0, first=True)
    # today marker: near the start of week 2 (Tue; Wed to Fri still ahead)
    marker_x = ML + (w + gap) + 0.25 * w
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 Inches(marker_x - 0.09), Inches(y + 0.62), Inches(0.18), Inches(0.14))
    tri.fill.solid()
    tri.fill.fore_color.rgb = CRIT
    tri.line.fill.background()
    tri.shadow.inherit = False
    tf = textbox(slide, marker_x - 0.5, y + 0.78, 1.0, 0.3)
    para(tf, L["s3_today"], size=10.5, bold=True, color=CRIT, align=PP_ALIGN.CENTER, space_after=0, first=True)

    # ---- 4 · architecture ----
    slide = content_slide(prs, L["s4_title"], L["footer"])
    intro(slide, L["s4_intro"], h=0.9)
    y, w, aw, gap = 3.0, 1.98, 0.38, 0.045
    x = ML
    for i, (title, sub) in enumerate(L["s4_nodes"]):
        flow_node(slide, x, y, w, 1.25, title, sub)
        x += w
        if i < len(L["s4_nodes"]) - 1:
            arrow(slide, x + gap, y + 0.50)
            x += aw + 2 * gap
    tf = textbox(slide, ML, 4.85, CW, 0.6)
    para(tf, L["s4_note"], size=12.5, color=INK, space_after=0, first=True)

    # ---- 5 · data access ----
    slide = content_slide(prs, L["s5_title"], L["footer"])
    intro(slide, L["s5_intro"], h=0.7)
    y = 2.35
    for name, desc, status, color, dark in L["s5_rows"]:
        shape = box(slide, ML, y, CW, 1.06, TINT)
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        para(tf, name, size=14.5, bold=True, color=NAVY, space_after=2, first=True)
        para(tf, desc, size=11, color=INK, space_after=0)
        pill(slide, ML + CW - 2.0, y + 0.14, 1.85, status, color, dark_text=dark)
        y += 1.2

    # ---- 6 · methodology ----
    slide = content_slide(prs, L["s6_title"], L["footer"])
    intro(slide, L["s6_intro"], h=0.95)
    zy, zh = 2.8, 0.55
    zw = [CW * 2 / 6, CW * 1 / 6, CW * 3 / 6]
    zc = [CRIT, WARN, GOOD]
    zdark = [False, True, False]
    x = ML
    for i, label in enumerate(L["s6_zones"]):
        shape = box(slide, x, zy, zw[i] - 0.03, zh, zc[i], radius=0.14)
        para(shape.text_frame, label, size=12.5, bold=True,
             color=INK if zdark[i] else PAPER, align=PP_ALIGN.CENTER, space_after=0, first=True)
        x += zw[i]
    tf = textbox(slide, ML + zw[0] - 0.6, zy + zh + 0.03, 1.2, 0.28)
    para(tf, L["s6_ticks"][0], size=9.5, color=INK, align=PP_ALIGN.CENTER, space_after=0, first=True)
    tf = textbox(slide, ML + zw[0] + zw[1] - 0.6, zy + zh + 0.03, 1.2, 0.28)
    para(tf, L["s6_ticks"][1], size=9.5, color=INK, align=PP_ALIGN.CENTER, space_after=0, first=True)
    tf = textbox(slide, ML, 3.85, CW, 0.35)
    para(tf, L["s6_weights_label"], size=12, bold=True, color=INK, space_after=0, first=True)
    wy, wh = 4.25, 0.5
    weights = [35, 20, 20, 10, 15]
    x = ML
    for i, weight in enumerate(weights):
        seg = CW * weight / 100
        shape = box(slide, x, wy, seg - 0.03, wh, W_COLORS[i], radius=0.14)
        para(shape.text_frame, str(weight), size=12, bold=True,
             color=INK if W_DARK_TEXT[i] else PAPER, align=PP_ALIGN.CENTER, space_after=0, first=True)
        tf = textbox(slide, x, wy + wh + 0.02, seg - 0.03, 0.28)
        para(tf, L["s6_weights"][i], size=9.5, color=INK, align=PP_ALIGN.CENTER, space_after=0, first=True)
        x += seg
    tf = textbox(slide, ML, 5.35, CW, 0.9)
    para(tf, L["s6_note"], size=12, color=INK, space_after=0, first=True)

    # ---- 7 · validation asks ----
    slide = content_slide(prs, L["s7_title"], L["footer"])
    intro(slide, L["s7_intro"], h=0.5)
    y, w, gap, h = 2.15, 5.83, 0.17, 1.75
    for i, (title, body) in enumerate(L["s7_cards"]):
        card(slide, ML + (i % 2) * (w + gap), y + (i // 2) * (h + 0.17), w, h, title, body, body_size=12)
    tf = textbox(slide, ML, 6.15, CW, 0.6)
    para(tf, L["s7_note"], size=12.5, color=INK, space_after=0, first=True)

    # ---- 8 · universe challenge ----
    slide = content_slide(prs, L["s8_title"], L["footer"])
    intro(slide, L["s8_intro"], h=0.9)
    tf = textbox(slide, ML, 2.62, CW, 0.32)
    para(tf, L["s8_now_label"], size=11.5, bold=True, color=INK, space_after=0, first=True)
    y, w, aw = 2.98, 2.55, 0.38
    x = ML
    for i, (title, sub) in enumerate(L["s8_now"]):
        flow_node(slide, x, y, w, 0.95, title, sub)
        x += w
        if i < len(L["s8_now"]) - 1:
            arrow(slide, x + 0.05, y + 0.35)
            x += aw + 0.1
    tf = textbox(slide, ML, 4.15, CW, 0.32)
    para(tf, L["s8_prop_label"], size=11.5, bold=True, color=INK, space_after=0, first=True)
    y, w = 4.51, 2.55
    x = ML
    for i, (title, sub) in enumerate(L["s8_prop"]):
        fill = NAVY if i == len(L["s8_prop"]) - 1 else TINT
        tcolor = PAPER if i == len(L["s8_prop"]) - 1 else NAVY
        shape = box(slide, x, y, w, 0.95, fill)
        tf2 = shape.text_frame
        para(tf2, title, size=12.5, bold=True, color=tcolor, align=PP_ALIGN.CENTER, space_after=2, first=True)
        para(tf2, sub, size=9, color=TINT if i == len(L["s8_prop"]) - 1 else INK,
             align=PP_ALIGN.CENTER, space_after=0)
        x += w
        if i < len(L["s8_prop"]) - 1:
            arrow(slide, x + 0.05, y + 0.35)
            x += aw + 0.1
    tf = textbox(slide, ML, 5.75, CW, 0.85)
    para(tf, L["s8_note"], size=12, color=INK, space_after=0, first=True)

    # ---- 9 · other lenses ----
    slide = content_slide(prs, L["s9_title"], L["footer"])
    intro(slide, L["s9_intro"], h=0.9)
    y, w, gap = 2.85, 3.83, 0.17
    for i, (title, body) in enumerate(L["s9_cards"]):
        card(slide, ML + i * (w + gap), y, w, 2.5, title, body, body_size=11.5)

    # ---- 10 · next steps ----
    slide = content_slide(prs, L["s10_title"], L["footer"])
    y, w, gap = 1.75, 3.83, 0.17
    for i, (title, body) in enumerate(L["s10_cols"]):
        card(slide, ML + i * (w + gap), y, w, 1.9, title, body, body_size=12)
    tf = textbox(slide, ML, 4.1, CW, 0.7)
    para(tf, L["s10_cadence"], size=13, color=INK, space_after=0, first=True)
    y, w, gap = 5.05, 1.35, 0.25
    for i in range(4):
        shape = box(slide, ML + i * (w + gap), y, w, 0.45, NAVY, radius=0.3)
        para(shape.text_frame, L["s10_chip"], size=10.5, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, space_after=0, first=True)

    out = HERE / f"midterm_presentation_{lang}.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    for lang in ("en", "pt"):
        for key, value in STR[lang].items():
            def flat(v):
                if isinstance(v, str):
                    yield v
                elif isinstance(v, (list, tuple)):
                    for item in v:
                        yield from flat(item)
            for text in flat(value):
                assert "—" not in text, f"em dash in {lang}:{key}"
                assert "•" not in text, f"bullet in {lang}:{key}"
        print("wrote", build(lang))
